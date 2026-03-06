"""Extract text from PowerPoint (.pptx) files."""

import io
from typing import BinaryIO, Union


def extract_text_from_pptx(source: Union[str, BinaryIO, bytes]) -> str:
    """
    Extract all text from a PowerPoint file (.pptx).
    source: file path (str), file-like object, or bytes.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx is required. Install with: pip install python-pptx")

    if isinstance(source, bytes):
        source = io.BytesIO(source)

    prs = Presentation(source)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    parts.append(text)
    return "\n\n".join(parts) if parts else ""
